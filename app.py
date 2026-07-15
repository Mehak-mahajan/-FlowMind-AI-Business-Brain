import streamlit as st
import pdfplumber
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from groq import Groq
from dotenv import load_dotenv
import os
import docx
import pptx
from langdetect import detect, detect_langs
import pandas as pd
from datetime import datetime
import pickle
import sqlite3
import json
import hashlib
import re
import io
import qrcode

load_dotenv()

# ---- PAGE SETUP ----
st.set_page_config(
    page_title="FlowMind — AI Business Brain",
    page_icon="🧠",
    layout="wide"
)

# ---- CSS ----
st.markdown("""
<style>
    .main { background-color: #ECE5DD; }
    .header-box {
        background-color: #075E54;
        padding: 15px 20px;
        border-radius: 10px;
        margin-bottom: 20px;
    }
    .user-bubble {
        background-color: #DCF8C6;
        padding: 10px 15px;
        border-radius: 15px 15px 0px 15px;
        margin: 5px 0;
        margin-left: 20%;
        color: #000;
        font-size: 14px;
        box-shadow: 0 1px 2px rgba(0,0,0,0.2);
    }
    .assistant-bubble {
        background-color: #FFFFFF;
        padding: 10px 15px;
        border-radius: 15px 15px 15px 0px;
        margin: 5px 0;
        margin-right: 20%;
        color: #000;
        font-size: 14px;
        box-shadow: 0 1px 2px rgba(0,0,0,0.2);
    }
    .source-tag {
        font-size: 11px;
        color: #075E54;
        margin-top: 3px;
    }
    .lead-box {
        background-color: #FFF3CD;
        padding: 15px;
        border-radius: 10px;
        border-left: 4px solid #FFC107;
        margin: 10px 0;
    }
    .alert-box {
        background-color: #FFE0E0;
        padding: 15px;
        border-radius: 10px;
        border-left: 4px solid #FF0000;
        margin: 10px 0;
    }
    .gap-box {
        background-color: #E0F0FF;
        padding: 12px;
        border-radius: 10px;
        border-left: 4px solid #2196F3;
        margin: 10px 0;
    }
    .welcome-box {
        text-align: center;
        padding: 40px;
        background: white;
        border-radius: 15px;
        margin: 20px 0;
    }
</style>
""", unsafe_allow_html=True)

# ---- LOAD MODELS ----
@st.cache_resource
def load_model():
    # multilingual-e5-small is built for asymmetric retrieval (short question
    # vs longer document text) — unlike paraphrase models, it's specifically
    # trained to match "WHEN WILL ARRIVE" against "Shipping: Pan India"
    # style content. Needs "query: " / "passage: " prefixes (added below)
    # to actually get this benefit — that's part of how the model was trained.
    return SentenceTransformer("intfloat/multilingual-e5-small")

embedding_model = load_model()
def get_api_key():
    # Works locally via .env, and on Streamlit Community Cloud via
    # the Secrets manager (Settings > Secrets in the deployed app).
    key = os.getenv("GROQ_API_KEY")
    if not key:
        try:
            key = st.secrets["GROQ_API_KEY"]
        except Exception:
            key = None
    return key

groq_client = Groq(api_key=get_api_key())

# Anchor storage to this script's own folder — not the current working
# directory, which can silently change depending on how streamlit gets
# launched (double-click, different terminal cwd, etc.) and make it look
# like accounts/data "disappeared" when really it's reading a different file.
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
KB_DIR = os.path.join(BASE_DIR, "kb_store")
os.makedirs(KB_DIR, exist_ok=True)

# =====================================================
# ---- SQLITE: MULTI-BUSINESS SCHEMA (NEW) ----
# =====================================================
DB_PATH = os.path.join(BASE_DIR, "flowmind.db")

def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS businesses (
            business_id TEXT PRIMARY KEY,
            name TEXT,
            password_hash TEXT,
            created_at TEXT
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS leads (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            business_id TEXT,
            name TEXT,
            phone TEXT,
            email TEXT,
            query TEXT,
            time TEXT
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS tickets (
            id TEXT PRIMARY KEY,
            business_id TEXT,
            issue TEXT,
            priority TEXT,
            time TEXT,
            status TEXT
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS unanswered_queries (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            business_id TEXT,
            query TEXT,
            time TEXT
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS products (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            business_id TEXT,
            name TEXT,
            price TEXT,
            description TEXT,
            created_at TEXT
        )
    """)
    conn.commit()

    # Migration: add follow_up_sent to leads if this is an older DB file
    try:
        c.execute("ALTER TABLE leads ADD COLUMN follow_up_sent INTEGER DEFAULT 0")
        conn.commit()
    except sqlite3.OperationalError:
        pass  # column already exists

    # Migration: private_notes on businesses — visible only to the owner,
    # never fed to the AI or shown to customers
    try:
        c.execute("ALTER TABLE businesses ADD COLUMN private_notes TEXT DEFAULT ''")
        conn.commit()
    except sqlite3.OperationalError:
        pass  # column already exists

    conn.close()

init_db()

# ---- BUSINESS HELPERS (NEW) ----
def slugify(name):
    slug = re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-")
    return slug or "business"

def hash_pw(pw):
    return hashlib.sha256(pw.encode()).hexdigest()

def normalize_business_id(business_id):
    # Business IDs are always generated lowercase — if someone types it back
    # with different casing or stray spaces (easy to do, it's not a password
    # they're used to typing carefully), lookups should still succeed.
    return (business_id or "").strip().lower()

def create_business(name, password):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    base_id = slugify(name)
    business_id = base_id
    n = 1
    # ensure unique id
    while c.execute(
        "SELECT 1 FROM businesses WHERE business_id=?", (business_id,)
    ).fetchone():
        n += 1
        business_id = f"{base_id}-{n}"
    c.execute(
        "INSERT INTO businesses (business_id, name, password_hash, created_at) VALUES (?, ?, ?, ?)",
        (business_id, name, hash_pw(password), datetime.now().strftime("%Y-%m-%d %H:%M"))
    )
    conn.commit()
    conn.close()
    return business_id

def verify_business(business_id, password):
    business_id = normalize_business_id(business_id)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    row = c.execute(
        "SELECT password_hash FROM businesses WHERE business_id=?", (business_id,)
    ).fetchone()
    conn.close()
    if row and row[0] == hash_pw(password):
        return True
    return False

def get_business_name(business_id):
    business_id = normalize_business_id(business_id)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    row = c.execute(
        "SELECT name FROM businesses WHERE business_id=?", (business_id,)
    ).fetchone()
    conn.close()
    return row[0] if row else business_id

def business_exists(business_id):
    business_id = normalize_business_id(business_id)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    row = c.execute(
        "SELECT 1 FROM businesses WHERE business_id=?", (business_id,)
    ).fetchone()
    conn.close()
    return row is not None

# ---- LEADS / TICKETS / GAPS (now business-scoped) ----
def save_lead_db(business_id, lead):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute(
        "INSERT INTO leads (business_id, name, phone, email, query, time) VALUES (?, ?, ?, ?, ?, ?)",
        (business_id, lead["name"], lead["phone"], lead["email"], lead["query"], lead["time"])
    )
    conn.commit()
    conn.close()

def get_leads_needing_followup(business_id, hours=24):
    """
    Leads older than `hours` that haven't been marked as followed up.
    This is the core of the follow-up feature: a lead that got no
    reply is money the business owner is quietly losing.
    """
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql_query(
        "SELECT * FROM leads WHERE business_id=? AND follow_up_sent=0 ORDER BY id ASC",
        conn, params=(business_id,)
    )
    conn.close()
    if df.empty:
        return df
    df["time_parsed"] = pd.to_datetime(df["time"])
    cutoff = pd.Timestamp.now() - pd.Timedelta(hours=hours)
    return df[df["time_parsed"] <= cutoff]

def mark_followed_up(lead_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE leads SET follow_up_sent=1 WHERE id=?", (lead_id,))
    conn.commit()
    conn.close()

def draft_followup_message(lead_query, business_name):
    """AI-drafted follow-up message the owner can send as-is or edit."""
    try:
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": f"""You write short, warm WhatsApp follow-up messages for
{business_name}, a small business. The customer asked something and showed buying
interest but never completed a purchase. Write ONE short, friendly follow-up message
(under 40 words) checking if they're still interested. No greeting like 'Dear Sir'.
Sound like a real shop owner texting, not a corporate email."""
                },
                {"role": "user", "content": f"Customer's original question: {lead_query}"}
            ],
            temperature=0.7,
            max_tokens=100
        )
        return response.choices[0].message.content.strip()
    except Exception:
        return "Hi! Just checking in — are you still interested in what you asked about earlier? Happy to help! 😊"

def generate_daily_digest(business_id, business_name):
    """
    One WhatsApp-style summary of the day for the owner —
    small business owners check WhatsApp, not dashboards.
    """
    today = datetime.now().strftime("%Y-%m-%d")
    leads_df = load_leads_db(business_id)
    tickets_df = load_tickets_db(business_id)
    gaps_df = load_gaps_db(business_id)

    today_leads = leads_df[leads_df["time"].str.startswith(today)] if not leads_df.empty else leads_df
    today_tickets = tickets_df[tickets_df["time"].str.startswith(today)] if not tickets_df.empty else tickets_df
    today_gaps = gaps_df[gaps_df["time"].str.startswith(today)] if not gaps_df.empty else gaps_df
    pending_followups = get_leads_needing_followup(business_id, hours=24)

    summary_input = f"""
Business: {business_name}
Today's leads: {len(today_leads)}
Today's support tickets: {len(today_tickets)}
Today's unanswered/gap questions: {len(today_gaps)}
Leads waiting 24h+ with no follow-up: {len(pending_followups)}
Sample gap questions: {', '.join(today_gaps['query'].head(3).tolist()) if not today_gaps.empty else 'none'}
"""
    try:
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": """Write a short, friendly end-of-day WhatsApp-style summary
for a small business owner based on the stats given. Use emojis sparingly, keep it under
80 words, plain language, no corporate tone. If there are leads needing follow-up or gap
questions, mention them clearly since those need the owner's attention."""
                },
                {"role": "user", "content": summary_input}
            ],
            temperature=0.5,
            max_tokens=200
        )
        return response.choices[0].message.content.strip()
    except Exception:
        return summary_input

def save_ticket_db(business_id, ticket):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute(
        "INSERT INTO tickets (id, business_id, issue, priority, time, status) VALUES (?, ?, ?, ?, ?, ?)",
        (ticket["id"], business_id, ticket["issue"], ticket["priority"], ticket["time"], ticket["status"])
    )
    conn.commit()
    conn.close()

def save_gap_db(business_id, query):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute(
        "INSERT INTO unanswered_queries (business_id, query, time) VALUES (?, ?, ?)",
        (business_id, query, datetime.now().strftime("%Y-%m-%d %H:%M"))
    )
    conn.commit()
    conn.close()

def load_leads_db(business_id):
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql_query(
        "SELECT * FROM leads WHERE business_id=? ORDER BY id DESC", conn, params=(business_id,)
    )
    conn.close()
    return df

def load_tickets_db(business_id):
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql_query(
        "SELECT * FROM tickets WHERE business_id=? ORDER BY id DESC", conn, params=(business_id,)
    )
    conn.close()
    return df

def load_gaps_db(business_id):
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql_query(
        "SELECT * FROM unanswered_queries WHERE business_id=? ORDER BY id DESC", conn, params=(business_id,)
    )
    conn.close()
    return df

def count_rows(table, business_id):
    conn = sqlite3.connect(DB_PATH)
    n = conn.execute(
        f"SELECT COUNT(*) FROM {table} WHERE business_id=?", (business_id,)
    ).fetchone()[0]
    conn.close()
    return n

# =====================================================
# ---- PRIVATE NOTES (NEW) — owner-only, never fed to the AI ----
# =====================================================
def save_private_notes(business_id, notes):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE businesses SET private_notes=? WHERE business_id=?", (notes, business_id))
    conn.commit()
    conn.close()

def get_private_notes(business_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    row = c.execute("SELECT private_notes FROM businesses WHERE business_id=?", (business_id,)).fetchone()
    conn.close()
    return row[0] if row and row[0] else ""

# =====================================================
# ---- PRODUCTS (NEW) — a browsable list customers see directly,
# not just something they have to ask the chatbot about ----
# =====================================================
def add_product(business_id, name, price, description):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute(
        "INSERT INTO products (business_id, name, price, description, created_at) VALUES (?, ?, ?, ?, ?)",
        (business_id, name, price, description, datetime.now().strftime("%Y-%m-%d %H:%M"))
    )
    conn.commit()
    conn.close()

def get_products(business_id):
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql_query(
        "SELECT * FROM products WHERE business_id=? ORDER BY id DESC", conn, params=(business_id,)
    )
    conn.close()
    return df

def delete_product(product_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("DELETE FROM products WHERE id=?", (product_id,))
    conn.commit()
    conn.close()

# =====================================================
# ---- PER-BUSINESS KNOWLEDGE BASE (NEW) ----
# =====================================================
def kb_paths(business_id):
    return (
        os.path.join(KB_DIR, f"{business_id}_chunks.pkl"),
        os.path.join(KB_DIR, f"{business_id}_index.bin")
    )

def save_knowledge_base(business_id, chunks, metadata, index):
    pkl_path, idx_path = kb_paths(business_id)
    with open(pkl_path, "wb") as f:
        pickle.dump({"chunks": chunks, "metadata": metadata}, f)
    faiss.write_index(index, idx_path)
    st.success("✅ Knowledge base saved for this business!")

def load_knowledge_base(business_id):
    pkl_path, idx_path = kb_paths(business_id)
    if os.path.exists(pkl_path) and os.path.exists(idx_path):
        with open(pkl_path, "rb") as f:
            data = pickle.load(f)
        index = faiss.read_index(idx_path)
        return data["chunks"], data["metadata"], index
    return None, None, None

# =====================================================
# ---- QR CODE (NEW) ----
# =====================================================
def generate_qr_code(url):
    """Customers scan this instead of typing a URL — same habit as UPI QR.
    Returns raw PNG bytes (not the qrcode library's internal image object) —
    that internal object is a 1-bit image which some Pillow/Streamlit
    version combinations fail to render directly."""
    qr = qrcode.QRCode(box_size=8, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white").convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

# =====================================================
# ---- QUICK SETUP (NEW) — no organized PDFs needed ----
# =====================================================
def build_business_profile_text(answers):
    lines = []
    for label, value in answers.items():
        if value and value.strip():
            lines.append(f"{label}: {value.strip()}")
    return "\n\n".join(lines)

def _normalize_vectors(vectors):
    vectors = np.array(vectors)
    norms = np.linalg.norm(vectors, axis=1, keepdims=True)
    return vectors / np.clip(norms, 1e-10, None)

def append_to_kb(business_id, new_chunks, new_metadata, new_vectors):
    """Adds to the existing knowledge base instead of replacing it —
    lets Quick Setup and file uploads build on each other over time."""
    existing_chunks, existing_metadata, existing_index = load_knowledge_base(business_id)
    if existing_chunks is None:
        dimension = new_vectors.shape[1]
        index = faiss.IndexFlatL2(dimension)
        index.add(new_vectors)
        save_knowledge_base(business_id, new_chunks, new_metadata, index)
    else:
        existing_index.add(new_vectors)
        save_knowledge_base(
            business_id,
            existing_chunks + new_chunks,
            existing_metadata + new_metadata,
            existing_index
        )

def process_quick_setup(business_id, answers):
    """Turns a simple Q&A form into knowledge base chunks — for business
    owners who don't have organized PDFs, which is most small businesses."""
    text_blob = build_business_profile_text(answers)
    if not text_blob.strip():
        return False

    chunk_size, overlap = 500, 100
    step = chunk_size - overlap
    chunks, metadata = [], []
    for i in range(0, len(text_blob), step):
        chunk = text_blob[i:i + chunk_size]
        if chunk.strip():
            chunks.append(chunk)
            metadata.append({"file": "Quick Setup Info", "page": 1})

    if not chunks:
        return False

    # e5 models need "passage: " prefix on document-side text to work correctly
    prefixed_chunks = [f"passage: {c}" for c in chunks]
    vectors = embedding_model.encode(prefixed_chunks, show_progress_bar=False, batch_size=16)
    vectors = _normalize_vectors(vectors)
    append_to_kb(business_id, chunks, metadata, vectors)
    return True

def embed_product_in_kb(business_id, name, price, description):
    """So the chatbot can also answer questions about a product,
    not just show it in the browsable list."""
    text = f"Product: {name} | Price: {price} | Details: {description}"
    vectors = embedding_model.encode([f"passage: {text}"], show_progress_bar=False)
    vectors = _normalize_vectors(vectors)
    append_to_kb(business_id, [text], [{"file": "Product Catalog", "page": 1}], vectors)

# ---- TEXT EXTRACTION ----
def extract_text(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()
    pages = []

    if file_type == "pdf":
        with pdfplumber.open(uploaded_file) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    pages.append((page_text, i + 1))

    elif file_type == "txt":
        text = uploaded_file.read().decode("utf-8")
        pages.append((text, 1))

    elif file_type == "docx":
        doc = docx.Document(uploaded_file)
        full_text = ""
        for paragraph in doc.paragraphs:
            full_text += paragraph.text + "\n"
        pages.append((full_text, 1))

    elif file_type == "pptx":
        prs = pptx.Presentation(uploaded_file)
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text + "\n"
            if slide_text:
                pages.append((slide_text, i + 1))

    elif file_type == "csv":
        text = uploaded_file.read().decode("utf-8")
        pages.append((text, 1))

    return pages

# ---- PROCESS FILES (chunking with overlap) ----
def process_files(uploaded_files):
    all_chunks = []
    all_metadata = []
    chunk_size = 500
    overlap = 100

    for uploaded_file in uploaded_files:
        pages = extract_text(uploaded_file)
        for page_text, page_num in pages:
            step = chunk_size - overlap
            for i in range(0, len(page_text), step):
                chunk = page_text[i:i + chunk_size]
                if chunk.strip():
                    all_chunks.append(chunk)
                    all_metadata.append({"file": uploaded_file.name, "page": page_num})
        st.write(f"✅ {uploaded_file.name} processed!")

    if not all_chunks:
        return None, None, None

    # e5 models need "passage: " prefix on document-side text to work correctly
    prefixed_chunks = [f"passage: {c}" for c in all_chunks]
    vectors = embedding_model.encode(prefixed_chunks, show_progress_bar=False, batch_size=32)
    vectors = _normalize_vectors(vectors)
    return all_chunks, all_metadata, vectors

# ---- LANGUAGE DETECTION (confidence-gated) ----
def safe_detect_language(text):
    """
    langdetect is a statistical model — it needs enough text to work
    reliably. Short phrases like 'is large size available' don't give
    it enough signal, and it can confidently guess the wrong language
    entirely. Default to English unless detection is both long enough
    to trust and confident enough to trust.
    """
    text = text.strip()
    if len(text) < 25 or len(text.split()) < 5:
        return "en"
    try:
        candidates = detect_langs(text)
        if candidates and candidates[0].prob >= 0.90:
            return candidates[0].lang
        return "en"
    except Exception:
        return "en"

# ---- LLM-BASED CLASSIFICATION ----
def classify_message(text):
    try:
        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": """You classify customer messages for an e-commerce support bot.
Return ONLY valid JSON, no markdown, no explanation, in this exact format:
{"sentiment": "angry" or "neutral", "is_lead": true or false}

sentiment = "angry" only for real frustration, anger, complaint, or strong dissatisfaction.
is_lead = true only for genuine buying intent (wants to purchase, bulk order, how to order).
A general or complaint-related price question is NOT a lead."""
                },
                {"role": "user", "content": text}
            ],
            temperature=0,
            max_tokens=50
        )
        raw = response.choices[0].message.content.strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        result = json.loads(raw)
        return result.get("sentiment", "neutral"), result.get("is_lead", False)
    except Exception:
        return "neutral", False

# =====================================================
# ---- CONFIDENCE-BASED RETRIEVAL (NEW) ----
# =====================================================
def retrieve_context(question, chunks, metadata, index, confidence_threshold):
    """
    Returns (context, sources, is_confident).
    is_confident=False means the best match is too weak to trust —
    caller should fall back instead of asking the LLM to guess.
    """
    # e5 models need "query: " prefix on question-side text to work correctly
    question_vector = embedding_model.encode([f"query: {question}"])
    question_vector = np.array(question_vector)
    question_vector = question_vector / np.clip(
        np.linalg.norm(question_vector, axis=1, keepdims=True), 1e-10, None
    )
    distances, indices = index.search(question_vector, 3)

    best_distance = float(distances[0][0])
    is_confident = best_distance <= confidence_threshold

    context = ""
    sources = []
    for i in indices[0]:
        context += chunks[i] + "\n\n"
        source = metadata[i]
        if source not in sources:
            sources.append(source)

    return context, sources, is_confident, best_distance

# =====================================================
# ---- HALLUCINATION GUARD (NEW) ----
# =====================================================
def response_has_unverified_contact_info(answer_text, context):
    """
    LLMs can invent specific-sounding details (a support URL, an email,
    a phone number) even when told to use only the given context — this
    catches that before it reaches the customer. A fabricated contact
    detail is worse than an honest 'let me connect you with support'.
    """
    patterns = [
        r'https?://[^\s]+',
        r'\b(?:www\.)?[a-zA-Z0-9][a-zA-Z0-9-]*\.(?:com|in|org|net|co|shop|store|info)\b',
        r'[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+',
        # Catches grouped/dashed formats too (e.g. "1800-123-4567"), not just
        # a bare run of 10 digits — the original version missed this format.
        r'(?<!\d)(?:\+?\d{1,3}[-\s]?)?\d(?:[-\s]?\d){7,}(?!\d)',
    ]
    context_lower = context.lower()
    for pattern in patterns:
        for match in re.findall(pattern, answer_text, flags=re.IGNORECASE):
            cleaned = match.lower().strip(".,;:")
            if cleaned and cleaned not in context_lower:
                return True
    return False

# ---- ASK QUESTION (with memory + confidence fallback) ----
def ask_question(question, context, lang, history):
    messages = [
        {
            "role": "system",
            "content": f"""You are FlowMind, an intelligent AI business assistant for an e-commerce store.
            Answer in same language as question (detected: {lang}).
            Use ONLY the provided context — never state a website, email, phone number,
            or any specific detail unless it is written word-for-word in the context.
            If contact details aren't in the context, say you'll connect them with support
            instead of guessing or inventing one.
            Use the recent conversation for follow-up questions (e.g. 'what about the blue one').
            Be helpful, friendly and concise."""
        }
    ]
    for turn in history[-6:]:
        # session_state messages carry extra fields (like "sources") for the UI —
        # Groq's API only accepts role/content, so strip anything else here.
        messages.append({"role": turn["role"], "content": turn["content"]})
    messages.append({"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}"})

    stream = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=messages,
        stream=True
    )
    for chunk in stream:
        content = chunk.choices[0].delta.content
        if content:
            yield content

# ---- SUMMARIZE / FAQ (unchanged logic) ----
def summarize_document(chunks):
    sample = "\n\n".join(chunks[:15])
    stream = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {"role": "system", "content": "Summarize this business document in clear bullet points."},
            {"role": "user", "content": f"Summarize:\n{sample}"}
        ],
        stream=True
    )
    for chunk in stream:
        content = chunk.choices[0].delta.content
        if content:
            yield content

def generate_faqs(chunks):
    sample = "\n\n".join(chunks[:15])
    stream = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {"role": "system", "content": "Generate top 10 FAQs with answers from this business document."},
            {"role": "user", "content": f"Generate FAQs:\n{sample}"}
        ],
        stream=True
    )
    for chunk in stream:
        content = chunk.choices[0].delta.content
        if content:
            yield content

# ---- SESSION STATE ----
for key, default in [
    ("messages", []), ("questions_count", 0), ("is_admin", False),
    ("business_id", None), ("confidence_threshold", 0.75)
]:
    if key not in st.session_state:
        st.session_state[key] = default

# ---- RESOLVE BUSINESS FROM URL (?biz=business-id) ----
query_biz = st.query_params.get("biz", None)
if query_biz and business_exists(query_biz):
    active_business_id = query_biz
elif st.session_state.is_admin and st.session_state.business_id:
    active_business_id = st.session_state.business_id
else:
    active_business_id = None

# ---- HEADER ----
header_name = get_business_name(active_business_id) if active_business_id else "FlowMind"
st.markdown(f"""
<div class="header-box">
    <h2 style="color:white;margin:0;">🧠 {header_name} — AI Business Brain</h2>
    <p style="color:#25D366;margin:0;">● Online | Powered by AI</p>
</div>
""", unsafe_allow_html=True)

# ---- SIDEBAR ----
# Customers arriving via a widget link (?biz=...) never see the admin panel at all —
# only the plain base URL (what only the business owner would open) shows it.
is_customer_link = query_biz is not None and business_exists(query_biz) and not st.session_state.is_admin

with st.sidebar:
    if is_customer_link:
        st.caption("💬 You're chatting as a customer. Business owners can manage their assistant from the main FlowMind link.")
    else:
        st.markdown("### 🔐 Admin Panel")

    if not st.session_state.is_admin and not is_customer_link:
        tab1, tab2 = st.tabs(["Login", "Create Business"])

        with tab1:
            login_id = st.text_input("Business ID", placeholder="e.g. shopeasy-india")
            login_pw = st.text_input("Password", type="password", key="login_pw")
            if st.button("Login"):
                if verify_business(login_id, login_pw):
                    st.session_state.is_admin = True
                    st.session_state.business_id = normalize_business_id(login_id)
                    st.rerun()
                else:
                    st.error("❌ Wrong business ID or password!")

        with tab2:
            new_name = st.text_input("Business Name", placeholder="e.g. ShopEasy India")
            new_pw = st.text_input("Set Password", type="password", key="new_pw")
            if st.button("Create Business Account"):
                if new_name and new_pw:
                    biz_id = create_business(new_name, new_pw)
                    st.session_state.is_admin = True
                    st.session_state.business_id = biz_id
                    st.success(f"✅ Created! Your Business ID: {biz_id}")
                    st.rerun()
                else:
                    st.warning("Enter a name and password.")

    if st.session_state.is_admin:
        business_id = st.session_state.business_id
        st.success(f"✅ Logged in as: {get_business_name(business_id)}")

        widget_path = f"?biz={business_id}"
        st.text_input("🔗 Shareable customer widget link (path)", value=widget_path)

        # NEW: QR code — customers scan instead of typing a URL. Familiar
        # behavior in India already (UPI QR codes), so no new habit to teach.
        st.caption("Enter your app's published URL once you deploy it, to get a scannable QR code:")
        public_base_url = st.text_input(
            "Your published app URL", placeholder="https://your-app.streamlit.app",
            key="public_base_url"
        )
        if public_base_url:
            full_url = public_base_url.rstrip("/") + "/" + widget_path
            qr_bytes = generate_qr_code(full_url)
            st.image(qr_bytes, caption="Print this on your counter, receipts, or shop window", width=200)
            st.download_button("📥 Download QR Code", qr_bytes, f"{business_id}_qr.png", "image/png")
        else:
            st.caption("Customers scan this to open your assistant directly — no typing needed.")

        st.divider()
        st.markdown("### 📁 Build Your Knowledge Base")
        tab_quick, tab_upload = st.tabs(["✍️ Quick Setup (no files needed)", "📁 Upload Documents"])

        with tab_quick:
            st.caption("Don't have organized PDFs? Just answer these — takes 2 minutes.")
            with st.form("quick_setup_form"):
                q_products = st.text_area("What do you sell? (products/services, briefly)")
                q_pricing = st.text_area("Pricing info (ranges are fine)")
                q_shipping = st.text_area("Shipping / delivery details")
                q_returns = st.text_area("Return / refund policy")
                q_hours = st.text_input("Business hours")
                q_contact = st.text_input("Contact info (phone/email/address)")
                q_faqs = st.text_area("Anything else customers often ask?")
                quick_submit = st.form_submit_button("➕ Add to Knowledge Base")

                if quick_submit:
                    answers = {
                        "Products/Services": q_products, "Pricing": q_pricing,
                        "Shipping": q_shipping, "Returns/Refunds": q_returns,
                        "Business Hours": q_hours, "Contact": q_contact, "Other FAQs": q_faqs
                    }
                    with st.spinner("Adding to knowledge base..."):
                        added = process_quick_setup(business_id, answers)
                    if added:
                        st.success("✅ Added! Your assistant can now answer from this.")
                        st.rerun()
                    else:
                        st.warning("Please fill in at least one field.")

        with tab_upload:
            uploaded_files = st.file_uploader(
                "Upload catalogs, policies, FAQs",
                type=["pdf", "txt", "docx", "pptx", "csv"],
                accept_multiple_files=True
            )
            st.caption("Adds to your existing knowledge base — works alongside Quick Setup, nothing gets overwritten.")
            if uploaded_files:
                if st.button("🧠 Process & Save Documents"):
                    with st.spinner("Processing..."):
                        new_chunks, new_metadata, new_vectors = process_files(uploaded_files)
                        if new_chunks:
                            append_to_kb(business_id, new_chunks, new_metadata, new_vectors)
                            st.rerun()

        st.divider()

        # NEW: Product Catalog — customers see this directly, not just through chat
        st.markdown("### 📦 Product Catalog")
        st.caption("Shown directly to customers above the chat — no need to ask the bot to discover what you sell.")
        with st.form("add_product_form"):
            p_name = st.text_input("Product name")
            p_price = st.text_input("Price (e.g. ₹1,200 or ₹500-₹2,000)")
            p_desc = st.text_area("Short description", height=68)
            add_product_submit = st.form_submit_button("➕ Add Product")
            if add_product_submit and p_name:
                add_product(business_id, p_name, p_price, p_desc)
                embed_product_in_kb(business_id, p_name, p_price, p_desc)
                st.success(f"✅ Added {p_name} — visible to customers and searchable by chat.")
                st.rerun()

        products_df = get_products(business_id)
        if not products_df.empty:
            for _, prod in products_df.iterrows():
                col1, col2 = st.columns([5, 1])
                with col1:
                    st.caption(f"**{prod['name']}** — {prod['price']}")
                with col2:
                    if st.button("🗑️", key=f"delprod_{prod['id']}"):
                        delete_product(prod["id"])
                        st.rerun()

        st.divider()

        # NEW: Private Notes — only the owner ever sees this, never fed to the AI
        st.markdown("### 🔒 Private Notes")
        st.caption("Only you can see this. Never shown to customers, never used by the AI. "
                   "Use it for internal info like supplier contacts, cost prices, or reminders.")
        current_notes = get_private_notes(business_id)
        new_notes = st.text_area("Your private notes", value=current_notes, height=100, key="private_notes_box")
        if st.button("💾 Save Private Notes"):
            save_private_notes(business_id, new_notes)
            st.success("Saved — visible only to you.")

        st.divider()

        with st.expander("⚠️ Danger Zone"):
            st.caption("Quick Setup and Document Upload both add to your knowledge base. "
                       "Use this only if you want to wipe everything and start over.")
            if st.button("🗑️ Reset Knowledge Base Completely"):
                pkl_path, idx_path = kb_paths(business_id)
                for p in (pkl_path, idx_path):
                    if os.path.exists(p):
                        os.remove(p)
                st.success("Knowledge base cleared. Add new info via Quick Setup or Upload.")
                st.rerun()

        st.divider()

        # NEW: confidence threshold control
        st.markdown("### 🎯 Answer Confidence")
        st.session_state.confidence_threshold = st.slider(
            "Lower = stricter (fewer guesses, more handoffs to support)",
            min_value=0.3, max_value=1.8,
            value=st.session_state.confidence_threshold, step=0.05
        )
        if "last_match_distance" in st.session_state:
            st.caption(f"Last question's match distance: {st.session_state.last_match_distance:.2f} "
                       f"({'✅ confident' if st.session_state.last_match_distance <= st.session_state.confidence_threshold else '⚠️ fell back'})")

        st.divider()
        st.markdown("### 📊 Analytics")
        leads_df = load_leads_db(business_id)
        tickets_df = load_tickets_db(business_id)
        gaps_df = load_gaps_db(business_id)

        col1, col2 = st.columns(2)
        with col1:
            st.metric("Questions Asked", st.session_state.questions_count)
            st.metric("Leads Captured", count_rows("leads", business_id))
        with col2:
            st.metric("Support Tickets", count_rows("tickets", business_id))
            st.metric("Unanswered Qs", count_rows("unanswered_queries", business_id))

        if not leads_df.empty:
            st.markdown("**Leads over time**")
            leads_df["date"] = pd.to_datetime(leads_df["time"]).dt.date
            st.bar_chart(leads_df.groupby("date").size())

        # NEW: content gap report — tells the business owner what's missing
        if not gaps_df.empty:
            st.markdown("### 🕳️ Knowledge Gaps")
            st.caption("Questions FlowMind couldn't confidently answer — add these to your docs.")
            st.dataframe(gaps_df[["query", "time"]], hide_index=True)

        st.divider()

        # NEW: Follow-up scheduler — the leads most businesses quietly lose money on
        st.markdown("### 🔔 Leads Needing Follow-up")
        st.caption("No reply after 24h+ — AI drafts the message, you just approve.")
        pending = get_leads_needing_followup(business_id, hours=24)
        if pending.empty:
            st.info("Nothing pending — you're all caught up! 🎉")
        else:
            for _, lead in pending.iterrows():
                with st.expander(f"{lead['name'] or 'Unnamed'} — asked: \"{lead['query'][:40]}...\""):
                    st.caption(f"Contact: {lead['phone'] or lead['email'] or 'not provided'} | {lead['time']}")
                    draft = draft_followup_message(lead["query"], get_business_name(business_id))
                    st.text_area("Suggested message", value=draft, key=f"draft_{lead['id']}", height=80)
                    colA, colB = st.columns(2)
                    with colA:
                        if st.button("✅ Mark as Sent", key=f"sent_{lead['id']}"):
                            mark_followed_up(lead["id"])
                            st.rerun()
                    with colB:
                        st.caption("Copy the text above to send via WhatsApp/SMS")

        st.divider()

        # NEW: Daily digest — the format a small business owner actually reads
        st.markdown("### 📋 Today's Digest")
        st.caption("A WhatsApp-style summary you could send yourself every evening.")
        if st.button("✨ Generate Today's Digest"):
            with st.spinner("Summarizing today..."):
                digest = generate_daily_digest(business_id, get_business_name(business_id))
                st.markdown(f'<div class="gap-box">{digest}</div>', unsafe_allow_html=True)

        st.divider()
        st.markdown("### 💰 ROI Calculator")
        staff_cost = st.number_input("Staff cost/month (₹)", value=25000, step=1000)
        st.success(f"💰 Monthly saving: ₹{staff_cost:,}")
        st.success(f"🚀 Yearly saving: ₹{staff_cost*12:,}")

        st.divider()
        st.markdown("### 🛠️ Admin Tools")
        chunks, metadata, index = load_knowledge_base(business_id)
        if chunks:
            if st.button("📝 Summarize Documents"):
                with st.expander("Summary", expanded=True):
                    st.write_stream(summarize_document(chunks))
            if st.button("❓ Generate FAQs"):
                with st.expander("FAQs", expanded=True):
                    st.write_stream(generate_faqs(chunks))

        if not leads_df.empty:
            st.download_button("📥 Download Leads CSV", leads_df.to_csv(index=False), "leads.csv", "text/csv")
        if not tickets_df.empty:
            st.download_button("📥 Download Tickets CSV", tickets_df.to_csv(index=False), "tickets.csv", "text/csv")

        if st.button("🗑️ Clear Chat History"):
            st.session_state.messages = []
            st.rerun()

        if st.button("🚪 Logout"):
            st.session_state.is_admin = False
            st.session_state.business_id = None
            st.rerun()

# ---- MAIN CHAT AREA ----
chunks, metadata, index = (None, None, None)
if active_business_id:
    chunks, metadata, index = load_knowledge_base(active_business_id)

if active_business_id and chunks:
    st.markdown(f"""
    <div style="background:#075E54;padding:10px;border-radius:8px;margin-bottom:10px;">
        <p style="color:white;margin:0;">🛒 {get_business_name(active_business_id)} — Customer Support</p>
        <p style="color:#25D366;margin:0;font-size:12px;">● Online | Ask us anything!</p>
    </div>
    """, unsafe_allow_html=True)

    # NEW: browsable product list — customers see this without having
    # to ask the chatbot to discover what's available
    customer_products = get_products(active_business_id)
    if not customer_products.empty:
        with st.expander(f"📦 Browse {get_business_name(active_business_id)}'s Products", expanded=False):
            for _, prod in customer_products.iterrows():
                st.markdown(f"**{prod['name']}** — {prod['price']}")
                if prod['description']:
                    st.caption(prod['description'])
                st.divider()

    for message in st.session_state.messages:
        if message["role"] == "user":
            st.markdown(f'<div class="user-bubble">👤 {message["content"]}</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div class="assistant-bubble">🧠 {message["content"]}</div>', unsafe_allow_html=True)
            for source in message.get("sources", []):
                st.markdown(f'<div class="source-tag">📄 {source["file"]} | Page {source["page"]}</div>', unsafe_allow_html=True)

    question = st.chat_input("Ask about products, shipping, returns...")

    if question:
        st.session_state.questions_count += 1
        st.markdown(f'<div class="user-bubble">👤 {question}</div>', unsafe_allow_html=True)
        st.session_state.messages.append({"role": "user", "content": question})

        sentiment, is_lead = classify_message(question)

        if sentiment == "angry":
            st.markdown("""
            <div class="alert-box">🚨 <strong>Priority Alert!</strong> Angry customer detected! Human agent needed!</div>
            """, unsafe_allow_html=True)
            ticket = {
                "id": f"TKT{count_rows('tickets', active_business_id)+1:03d}",
                "issue": question, "priority": "HIGH",
                "time": datetime.now().strftime("%Y-%m-%d %H:%M"), "status": "Open"
            }
            save_ticket_db(active_business_id, ticket)
            st.error(f"🎫 Ticket: {ticket['id']} created!")

        if is_lead:
            st.markdown("""
            <div class="lead-box">🎯 <strong>Potential Lead Detected!</strong> Customer interested in buying!</div>
            """, unsafe_allow_html=True)
            with st.form("lead_form"):
                st.markdown("**📋 Save Customer Details:**")
                name = st.text_input("Name")
                phone = st.text_input("Phone")
                email = st.text_input("Email")
                submit = st.form_submit_button("Save Lead")
                if submit and name:
                    lead = {"name": name, "phone": phone, "email": email,
                            "query": question, "time": datetime.now().strftime("%Y-%m-%d %H:%M")}
                    save_lead_db(active_business_id, lead)
                    st.success(f"✅ Lead saved! Total: {count_rows('leads', active_business_id)}")

        lang = safe_detect_language(question)

        context, sources, is_confident, best_distance = retrieve_context(
            question, chunks, metadata, index, st.session_state.confidence_threshold
        )
        st.session_state.last_match_distance = best_distance

        answer_text = ""
        placeholder = st.empty()

        if not is_confident:
            # NEW: confidence-based fallback instead of guessing
            answer_text = "I'm not fully sure about this from our current info — let me connect you with our support team so they can help directly. 🙏"
            placeholder.markdown(f'<div class="assistant-bubble">🧠 {answer_text}</div>', unsafe_allow_html=True)
            save_gap_db(active_business_id, question)
            sources = []
        else:
            for content in ask_question(question, context, lang, st.session_state.messages):
                answer_text += content
                placeholder.markdown(f'<div class="assistant-bubble">🧠 {answer_text}</div>', unsafe_allow_html=True)

            # NEW: catch invented contact details (fake URLs/emails/phone numbers)
            # after the full answer is in — a wrong website is worse than "let me check"
            if response_has_unverified_contact_info(answer_text, context):
                answer_text = ("I want to make sure I give you the correct contact details — "
                                "let me connect you with our support team directly for this. 🙏")
                placeholder.markdown(f'<div class="assistant-bubble">🧠 {answer_text}</div>', unsafe_allow_html=True)
                save_gap_db(active_business_id, question)
                sources = []
            else:
                for source in sources:
                    st.markdown(f'<div class="source-tag">📄 {source["file"]} | Page {source["page"]}</div>', unsafe_allow_html=True)

        st.session_state.messages.append({"role": "assistant", "content": answer_text, "sources": sources})

elif active_business_id and not chunks:
    st.info(f"👋 Welcome! {get_business_name(active_business_id)} hasn't uploaded documents yet. Please check back soon.")

else:
    st.markdown("""
    <div class="welcome-box">
        <h1>🧠 FlowMind</h1>
        <h3 style="color:#075E54;">AI Business Brain for Any Business</h3>
        <br>
        <p style="font-size:16px;">One AI assistant platform, unlimited businesses.</p>
        <p style="font-size:14px;color:#888;">
            Business owner? Create your account from the sidebar to get your own
            AI support assistant and a shareable widget link in under a minute.
        </p>
    </div>
    """, unsafe_allow_html=True)
    st.info("👈 Create a business account or login from the sidebar to get started!")